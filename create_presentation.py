import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_deck():
    prs = Presentation()
    # Use widescreen 16:9 aspect ratio
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Color Palette Definitions
    COLOR_BG_DARK = RGBColor(11, 19, 43)       # Deep Obsidian Blue
    COLOR_BG_LIGHT = RGBColor(248, 250, 252)   # Sleek Off-White (Slate 50)
    COLOR_PRIMARY = RGBColor(79, 70, 229)      # Vivid Royal Indigo (Indigo 600)
    COLOR_SECONDARY = RGBColor(99, 102, 241)    # Bright Indigo (Indigo 500)
    COLOR_ACCENT_GREEN = RGBColor(16, 185, 129)# Emerald Green (success/growth)
    COLOR_ACCENT_PURPLE = RGBColor(168, 85, 247)# Vibrant Purple
    
    COLOR_TEXT_WHITE = RGBColor(255, 255, 255)
    COLOR_TEXT_DARK = RGBColor(15, 23, 42)     # Slate 900
    COLOR_TEXT_MUTED = RGBColor(100, 116, 139) # Slate 500
    COLOR_CARD_BORDER = RGBColor(226, 232, 240)# Slate 200

    def apply_solid_bg(slide, rgb_color):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = rgb_color

    def add_title(slide, text, color=COLOR_PRIMARY, size=32, top=0.6, left=0.8, width=11.7, height=0.8):
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0)
        tf.margin_left = Inches(0)
        
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = "Segoe UI"
        p.font.size = Pt(size)
        p.font.bold = True
        p.font.color.rgb = color
        return txBox

    def add_bullet_list(slide, items, left, top, width, height, font_size=14, is_dark=False):
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0)
        tf.margin_top = Inches(0)

        for idx, item in enumerate(items):
            if idx == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = "•  " + item
            p.font.name = "Segoe UI"
            p.font.size = Pt(font_size)
            p.font.color.rgb = COLOR_TEXT_WHITE if is_dark else COLOR_TEXT_DARK
            p.space_after = Pt(10)
            p.line_spacing = 1.25

    def add_card(slide, title, description, left, top, width, height, bg_color=RGBColor(255, 255, 255), is_dark=False, border_color=COLOR_CARD_BORDER, badge_text=None, badge_color=COLOR_PRIMARY):
        # Draw rounded rectangle card shape
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, 
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        if is_dark:
            card.line.fill.background()
        else:
            card.line.color.rgb = border_color
            card.line.width = Pt(1.5)

        # Draw pill badge if requested
        if badge_text:
            badge = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(left + 0.3), Inches(top + 0.3), Inches(1.8), Inches(0.32)
            )
            badge.fill.solid()
            badge.fill.fore_color.rgb = badge_color
            badge.line.fill.background()
            
            btf = badge.text_frame
            btf.margin_top = Inches(0.02)
            bp = btf.paragraphs[0]
            bp.alignment = PP_ALIGN.CENTER
            bp.text = badge_text.upper()
            bp.font.name = "Segoe UI"
            bp.font.size = Pt(9)
            bp.font.bold = True
            bp.font.color.rgb = COLOR_TEXT_WHITE
            
            top_offset = 0.8
        else:
            top_offset = 0.3
            
        # Add content text box nested inside
        txBox = slide.shapes.add_textbox(Inches(left + 0.3), Inches(top + top_offset), Inches(width - 0.6), Inches(height - top_offset - 0.2))
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0)
        tf.margin_left = Inches(0)
        
        # Title
        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.name = "Segoe UI"
        p1.font.size = Pt(18)
        p1.font.bold = True
        p1.font.color.rgb = COLOR_TEXT_WHITE if is_dark else COLOR_PRIMARY
        p1.space_after = Pt(8)
        
        # Description
        p2 = tf.add_paragraph()
        p2.text = description
        p2.font.name = "Segoe UI"
        p2.font.size = Pt(13)
        p2.font.color.rgb = RGBColor(226, 232, 240) if is_dark else COLOR_TEXT_DARK
        p2.line_spacing = 1.25
        return card

    # ==========================================
    # SLIDE 1: Title Slide (Obsidian Dark Theme)
    # ==========================================
    slide_layout = prs.slide_layouts[6] # Blank
    slide1 = prs.slides.add_slide(slide_layout)
    apply_solid_bg(slide1, COLOR_BG_DARK)

    # Decorative background shapes for aesthetic depth
    shape1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(0.15))
    shape1.fill.solid()
    shape1.fill.fore_color.rgb = COLOR_PRIMARY
    shape1.line.fill.background()

    # Main Title
    title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.33), Inches(3.5))
    tf1 = title_box.text_frame
    tf1.word_wrap = True
    
    p = tf1.paragraphs[0]
    p.text = "SMART OFFICE"
    p.font.name = "Segoe UI"
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = COLOR_SECONDARY
    p.space_after = Pt(8)
    
    p2 = tf1.add_paragraph()
    p2.text = "The Branded Operating System for Coaching Institutes"
    p2.font.name = "Segoe UI"
    p2.font.size = Pt(24)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_TEXT_WHITE
    p2.space_after = Pt(16)

    p3 = tf1.add_paragraph()
    p3.text = "Empower Staff • Secure Revenue Ledgers • Track Inventory • Automate Payroll"
    p3.font.name = "Segoe UI"
    p3.font.size = Pt(16)
    p3.font.color.rgb = COLOR_TEXT_MUTED
    p3.space_after = Pt(20)

    p4 = tf1.add_paragraph()
    p4.text = "DELIVERING EFFICIENCY • ACCURACY • PROFESSIONALISM"
    p4.font.name = "Segoe UI"
    p4.font.size = Pt(12)
    p4.font.bold = True
    p4.font.color.rgb = COLOR_ACCENT_GREEN

    # ==========================================
    # SLIDE 2: The Operational Nightmare (Problem)
    # ==========================================
    slide2 = prs.slides.add_slide(slide_layout)
    apply_solid_bg(slide2, COLOR_BG_LIGHT)
    add_title(slide2, "Why Traditional Coaching Centers Struggle", COLOR_TEXT_DARK)
    
    # 3 Column Cards
    add_card(slide2, "Spreadsheet & Roster Chaos", 
             "Managing staff rosters, class schedules, and attendance sheets across separate, offline spreadsheets leads to constant synchronization mismatches, duplicate data entry, and payroll disputes.", 
             1.0, 1.8, 3.5, 4.6, badge_text="Administrative Drain", badge_color=COLOR_PRIMARY)
             
    add_card(slide2, "Revenue & Fee Leakage", 
             "Without standardized billing, tracking student enrollments, course fee structures, class duration rates, and pending student ledgers becomes a disorganized manual effort prone to revenue loss.", 
             4.9, 1.8, 3.5, 4.6, badge_text="Financial Leakage", badge_color=COLOR_PRIMARY)
             
    add_card(slide2, "Manual Monthly Accounting", 
             "Hours are wasted at the end of each month calculating basic pay, deductions for teacher absences, utility bills (electricity, water), and inventory expenses on paper ledgers.", 
             8.8, 1.8, 3.5, 4.6, badge_text="Manual Labor", badge_color=COLOR_PRIMARY)

    # ==========================================
    # SLIDE 3: White-Label Branding (Benefit-Focused)
    # ==========================================
    slide3 = prs.slides.add_slide(slide_layout)
    apply_solid_bg(slide3, COLOR_BG_LIGHT)
    add_title(slide3, "Your Custom Branded Portal", COLOR_PRIMARY)

    # Large feature card
    add_card(slide3, "Professional Custom Branding Out-of-the-Box",
             "Smart Office represents your coaching center's identity. The portal automatically adapts to display your specific center details to your staff and faculty, building trust and institutional authority.\n\n"
             "• Dynamic Header Logos: Custom school crests or corporate logos loaded immediately.\n"
             "• Brand Matching UI: The interface dynamically matches your school's color scheme.\n"
             "• Dynamic Sidebar & Slips: Your unique name, contact numbers, and office addresses are automatically hardcoded onto all pages, printable payslips, and ledgers.",
             0.8, 1.8, 6.0, 4.8, border_color=COLOR_PRIMARY)

    # Visual representation of rebranding on the right
    repr_card = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.4), Inches(1.8), Inches(5.1), Inches(4.8))
    repr_card.fill.solid()
    repr_card.fill.fore_color.rgb = COLOR_BG_DARK
    repr_card.line.fill.background()

    txBox = slide3.shapes.add_textbox(Inches(7.7), Inches(2.1), Inches(4.5), Inches(4.2))
    tf3 = txBox.text_frame
    tf3.word_wrap = True
    
    p = tf3.paragraphs[0]
    p.text = "INSTITUTION BRAND METRICS"
    p.font.name = "Segoe UI"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT_GREEN
    p.space_after = Pt(20)

    details = [
        "Brand Name: Loaded Dynamically",
        "Theme Color: Locked to Brand Accent",
        "Logo Upload: Super Admin Set",
        "Document Headers: Branded Automatically",
        "Staff Slips: Standardized Addresses"
    ]
    for d in details:
        pr = tf3.add_paragraph()
        pr.text = "✔  " + d
        pr.font.name = "Segoe UI"
        pr.font.size = Pt(14)
        pr.font.color.rgb = COLOR_TEXT_WHITE
        pr.space_after = Pt(14)

    # ==========================================
    # SLIDE 4: Staff & Automated Payroll (Benefit)
    # ==========================================
    slide4 = prs.slides.add_slide(slide_layout)
    apply_solid_bg(slide4, COLOR_BG_LIGHT)
    add_title(slide4, "Automated Roster & Payroll Processing", COLOR_PRIMARY)

    # Two columns
    add_card(slide4, "Real-Time Attendance Sheets", 
             "Ditch physical sign-in books. Track daily check-in statuses (Present, Absent, Late, Half-Day) and remarks. Staff and faculty directories calculate active workdays on-the-fly, reducing attendance disputes to zero.", 
             1.0, 1.8, 5.2, 4.6, badge_text="Staff Directory", badge_color=COLOR_ACCENT_PURPLE)

    add_card(slide4, "Zero-Error Payroll Calculations", 
             "Payroll cycles automatically pull total absences and compute prorated absent deductions, salary allowances, and net payouts instantly. Process salaries and print structured, professional payslips for faculty and clerks in one click.", 
             7.1, 1.8, 5.2, 4.6, badge_text="One-Click Payslips", badge_color=COLOR_ACCENT_GREEN)

    # ==========================================
    # SLIDE 5: Course & Student Ledgers (Benefit)
    # ==========================================
    slide5 = prs.slides.add_slide(slide_layout)
    apply_solid_bg(slide5, COLOR_BG_LIGHT)
    add_title(slide5, "Standardized Course Fees & Student Records", COLOR_PRIMARY)

    # Two columns
    add_card(slide5, "Course Fee Master Control", 
             "Standardize your core curriculum rates. Configure standard course prices, class durations, and billing options. This ensures clerks enroll students only under approved pricing tiers, stopping administrative discount overrides.", 
             1.0, 1.8, 5.2, 4.6, badge_text="Course Master", badge_color=COLOR_PRIMARY)

    add_card(slide5, "Chronological Student Ledgers", 
             "Every enrolled student has a secure profile with billing and fee payment statements. Track chronological payment history and student roll-calls. Easily verify due payments and prevent billing leakages.", 
             7.1, 1.8, 5.2, 4.6, badge_text="Revenue Protection", badge_color=COLOR_ACCENT_GREEN)

    # ==========================================
    # SLIDE 6: Monthly Accounts & Ledger (Benefit)
    # ==========================================
    slide6 = prs.slides.add_slide(slide_layout)
    apply_solid_bg(slide6, COLOR_BG_LIGHT)
    add_title(slide6, "Accounts, Inventory, & P&L Ledgers", COLOR_PRIMARY)

    # 3 Column Cards
    add_card(slide6, "Utility & Expense Logging", 
             "Log rent, electricity, water, and miscellaneous purchases. Upload expense descriptions immediately to maintain complete bookkeeping transparency.", 
             1.0, 1.8, 3.5, 4.6, badge_text="Expense Tracking", badge_color=COLOR_PRIMARY)
             
    add_card(slide6, "Profit & Loss Ledger", 
             "Automated monthly accounting. Compare total student fees collected (Revenue) against payroll payouts and utility expenses (Cost) to show net profitability instantly.", 
             4.9, 1.8, 3.5, 4.6, badge_text="Monthly P&L", badge_color=COLOR_ACCENT_GREEN)
             
    add_card(slide6, "Inventory Asset Control", 
             "Track coaching supplies, furniture, and textbook stocks. Manage asset item counts, unit costs, and receive stock warnings to prevent classroom supply shortages.", 
             8.8, 1.8, 3.5, 4.6, badge_text="Asset Audit", badge_color=COLOR_PRIMARY)

    # ==========================================
    # SLIDE 7: Standalone PWA PWA (Benefit)
    # ==========================================
    slide7 = prs.slides.add_slide(slide_layout)
    apply_solid_bg(slide7, COLOR_BG_LIGHT)
    add_title(slide7, "PWA Standalone App Experience", COLOR_PRIMARY)

    # 3 Column Cards on PWA Benefits
    add_card(slide7, "Browser-to-App Install", 
             "No App Store search or downloads required. Tap the PWA prompt directly in the browser to place the Smart Office icon on your Android, iOS, or Desktop home screen.", 
             1.0, 1.8, 3.5, 4.6, badge_text="Instant Access", badge_color=COLOR_SECONDARY)
             
    add_card(slide7, "Native Mobile Shell", 
             "Launches with branded splash screens, dynamic home icons, and runs without browser URL bars or navigation headers, providing a 100% native feel.", 
             4.9, 1.8, 3.5, 4.6, badge_text="Fluid Interface", badge_color=COLOR_SECONDARY)
             
    add_card(slide7, "Advanced Cache & Offline", 
             "Service workers cache app files locally, enabling administrators to check attendance, view staff logs, and input ledger data even in zero-internet environments.", 
             8.8, 1.8, 3.5, 4.6, badge_text="Offline Capable", badge_color=COLOR_ACCENT_GREEN)

    # ==========================================
    # SLIDE 8: Summary of Value (Business ROI)
    # ==========================================
    slide8 = prs.slides.add_slide(slide_layout)
    apply_solid_bg(slide8, COLOR_BG_LIGHT)
    add_title(slide8, "The Smart Office ROI: Operational Efficiency", COLOR_PRIMARY)

    # Left Description Column
    add_card(slide8, "How Smart Office Transforms Your Business", 
             "By unifying your staff rosters, student billing, inventory, and accounts into a single, white-labeled database, you eliminate human calculation errors, secure your financial records, and free up critical administrative hours.\n\n"
             "The system requires zero local database setups and features real-time Firestore database replication, ensuring your coaching center operations run smoothly and securely.", 
             0.8, 1.8, 5.5, 4.8, border_color=COLOR_PRIMARY)

    # Right Highlights List
    roi_metrics = [
        "Save 40+ hours per month on manual payroll, tracking absences, and bookkeeping.",
        "Zero-error payroll: automated deductions stop salary overpayments.",
        "Revenue protection: Course Fee Master prevents unauthorized student discounts.",
        "Total transparency: dynamic P&L ledger tracking utility bills and income logs.",
        "Offline productivity: clerks log attendance and view ledgers without active internet."
    ]
    add_bullet_list(slide8, roi_metrics, 6.7, 1.9, 5.8, 4.5, font_size=15.5)

    # ==========================================
    # SLIDE 9: Conclusion & Call to Action (Dark Theme)
    # ==========================================
    slide9 = prs.slides.add_slide(slide_layout)
    apply_solid_bg(slide9, COLOR_BG_DARK)

    # Center box
    accent2 = slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(0.15))
    accent2.fill.solid()
    accent2.fill.fore_color.rgb = COLOR_PRIMARY
    accent2.line.fill.background()

    cta_box = slide9.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.33), Inches(4.0))
    tf9 = cta_box.text_frame
    tf9.word_wrap = True
    
    p = tf9.paragraphs[0]
    p.text = "Bring Smart Office to Your Coaching Center"
    p.font.name = "Segoe UI"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLOR_SECONDARY
    p.space_after = Pt(12)
    
    p2 = tf9.add_paragraph()
    p2.text = "Streamline your institute's daily operations, automate your billing ledger, and manage payroll with zero mathematical errors in a clean, branded portal."
    p2.font.name = "Segoe UI"
    p2.font.size = Pt(18)
    p2.font.color.rgb = COLOR_TEXT_WHITE
    p2.space_after = Pt(24)

    p3 = tf9.add_paragraph()
    p3.text = "Contact the System Administrator to provision your branded portal and PWA app installer."
    p3.font.name = "Segoe UI"
    p3.font.size = Pt(16)
    p3.font.bold = True
    p3.font.color.rgb = COLOR_ACCENT_GREEN

    # Save presentation
    filepath = "C:/Users/deepe/.gemini/antigravity/scratch/Smart_Office_Coaching_Benefits.pptx"
    prs.save(filepath)
    print(f"Presentation created successfully at: {os.path.abspath(filepath)}")
    
    # Copy to Desktop
    import shutil
    desktop_path = "C:/Users/deepe/Desktop/Smart_Office_Coaching_Benefits.pptx"
    try:
        shutil.copy(filepath, desktop_path)
        print(f"Copied presentation to Desktop: {desktop_path}")
    except Exception as e:
        print(f"Failed to copy to Desktop: {e}")

if __name__ == "__main__":
    create_deck()
