import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_deck_with_images():
    prs = Presentation()
    # Use widescreen 16:9 aspect ratio
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Color Palette Definitions
    COLOR_BG_DARK = RGBColor(11, 19, 43)       # Deep Obsidian Blue
    COLOR_BG_LIGHT = RGBColor(248, 250, 252)   # Sleek Slate 50
    COLOR_PRIMARY = RGBColor(79, 70, 229)      # Vivid Royal Indigo
    COLOR_SECONDARY = RGBColor(99, 102, 241)    # Bright Indigo
    COLOR_ACCENT_GREEN = RGBColor(16, 185, 129)# Emerald Green
    COLOR_ACCENT_PURPLE = RGBColor(168, 85, 247)# Purple
    COLOR_TEXT_WHITE = RGBColor(255, 255, 255)
    COLOR_TEXT_DARK = RGBColor(15, 23, 42)     # Slate 900
    COLOR_TEXT_MUTED = RGBColor(100, 116, 139) # Slate 500
    COLOR_CARD_BORDER = RGBColor(226, 232, 240)# Slate 200

    brain_dir = "C:/Users/deepe/.gemini/antigravity/brain/c7c39514-81b3-45e7-9f4b-1afe886fe7a5"
    images = {
        "dashboard": os.path.join(brain_dir, "media__1781180982241.png"),
        "attendance": os.path.join(brain_dir, "media__1781181013900.png"),
        "payroll": os.path.join(brain_dir, "media__1781181024742.png"),
        "enrollment": os.path.join(brain_dir, "media__1781181038523.png"),
        "ledger": os.path.join(brain_dir, "media__1781181200179.png")
    }

    def apply_solid_bg(slide, rgb_color):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = rgb_color

    def add_title(slide, text, color=COLOR_PRIMARY, size=30, top=0.5, left=0.6, width=12.0, height=0.7):
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0)
        tf.margin_left = Inches(0)
        
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = "Segoe UI"
        p.font.size = Pt(size)
        p.font.bold = True
        p.font.color.rgb = color
        return txBox

    def add_bullet_list(slide, items, left, top, width, height, font_size=13.5, is_dark=False):
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0)
        tf.margin_top = Inches(0)

        for idx, item in enumerate(items):
            if idx == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = "•  " + item
            p.font.name = "Segoe UI"
            p.font.size = Pt(font_size)
            p.font.color.rgb = COLOR_TEXT_WHITE if is_dark else COLOR_TEXT_DARK
            p.space_after = Pt(8)
            p.line_spacing = 1.2

    def add_card(slide, title, description, left, top, width, height, bg_color=RGBColor(255, 255, 255), is_dark=False, border_color=COLOR_CARD_BORDER, badge_text=None, badge_color=COLOR_PRIMARY):
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, 
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        if is_dark:
            card.line.fill.background()
        else:
            card.line.color.rgb = border_color
            card.line.width = Pt(1.5)

        if badge_text:
            badge = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(left + 0.25), Inches(top + 0.25), Inches(1.8), Inches(0.3)
            )
            badge.fill.solid()
            badge.fill.fore_color.rgb = badge_color
            badge.line.fill.background()
            
            btf = badge.text_frame
            btf.margin_top = Inches(0.01)
            bp = btf.paragraphs[0]
            bp.alignment = PP_ALIGN.CENTER
            bp.text = badge_text.upper()
            bp.font.name = "Segoe UI"
            bp.font.size = Pt(8.5)
            bp.font.bold = True
            bp.font.color.rgb = COLOR_TEXT_WHITE
            
            top_offset = 0.7
        else:
            top_offset = 0.25
            
        txBox = slide.shapes.add_textbox(Inches(left + 0.25), Inches(top + top_offset), Inches(width - 0.5), Inches(height - top_offset - 0.15))
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0)
        tf.margin_left = Inches(0)
        
        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.name = "Segoe UI"
        p1.font.size = Pt(17)
        p1.font.bold = True
        p1.font.color.rgb = COLOR_TEXT_WHITE if is_dark else COLOR_PRIMARY
        p1.space_after = Pt(6)
        
        p2 = tf.add_paragraph()
        p2.text = description
        p2.font.name = "Segoe UI"
        p2.font.size = Pt(12.5)
        p2.font.color.rgb = RGBColor(226, 232, 240) if is_dark else COLOR_TEXT_DARK
        p2.line_spacing = 1.2
        return card

    def add_screenshot(slide, img_path, left, top, width, height):
        if os.path.exists(img_path):
            # Add subtle frame background behind image
            frame = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, 
                Inches(left - 0.04), Inches(top - 0.04), Inches(width + 0.08), Inches(height + 0.08)
            )
            frame.fill.solid()
            frame.fill.fore_color.rgb = COLOR_SECONDARY
            frame.line.fill.background()
            
            # Place image
            slide.shapes.add_picture(img_path, Inches(left), Inches(top), Inches(width), Inches(height))
        else:
            # Fallback box if image missing
            box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
            box.fill.solid()
            box.fill.fore_color.rgb = COLOR_TEXT_MUTED
            box.text = f"Screenshot Placeholder:\n{os.path.basename(img_path)}"
            print(f"Warning: Image file not found: {img_path}")

    # ==========================================
    # SLIDE 1: Cover Page (Obsidian Dark Theme)
    # ==========================================
    slide_layout = prs.slide_layouts[6]
    slide1 = prs.slides.add_slide(slide_layout)
    apply_solid_bg(slide1, COLOR_BG_DARK)

    # Accent Top Line
    accent = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(0.15))
    accent.fill.solid()
    accent.fill.fore_color.rgb = COLOR_PRIMARY
    accent.line.fill.background()

    # Title texts
    title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.33), Inches(4.0))
    tf1 = title_box.text_frame
    tf1.word_wrap = True
    
    p = tf1.paragraphs[0]
    p.text = "SMART OFFICE"
    p.font.name = "Segoe UI"
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = COLOR_SECONDARY
    p.space_after = Pt(8)
    
    p2 = tf1.add_paragraph()
    p2.text = "Branded Operations & Administration Platform"
    p2.font.name = "Segoe UI"
    p2.font.size = Pt(22)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_TEXT_WHITE
    p2.space_after = Pt(12)

    p3 = tf1.add_paragraph()
    p3.text = "Staff Management • Real-Time Attendance • Automated Payroll • Course Fee Master • Financial Ledgers"
    p3.font.name = "Segoe UI"
    p3.font.size = Pt(15)
    p3.font.color.rgb = COLOR_TEXT_MUTED
    p3.space_after = Pt(24)

    p4 = tf1.add_paragraph()
    p4.text = "COACHING INSTITUTE MARKETING DECK • EMBEDDED SOFTWARE DEMONSTRATION"
    p4.font.name = "Segoe UI"
    p4.font.size = Pt(11)
    p4.font.bold = True
    p4.font.color.rgb = COLOR_ACCENT_GREEN

    # ==========================================
    # SLIDE 2: Dashboard Overview (Light Theme)
    # ==========================================
    slide2 = prs.slides.add_slide(slide_layout)
    apply_solid_bg(slide2, COLOR_BG_LIGHT)
    add_title(slide2, "Unified Dashboard Overview", COLOR_PRIMARY)

    # Content on Left
    add_card(slide2, "Administrative Control Panel", 
             "The landing dashboard provides owners and administrators with immediate operational intelligence. It brings staff, student attendance, financial estimates, and system logs under one screen.",
             0.6, 1.4, 5.2, 2.5, badge_text="Live Analytics", badge_color=COLOR_SECONDARY)

    bullets2 = [
        "Real-Time Statistics: Monitor total active staff members, present rates, and current absentees.",
        "Students Overview: Instant counts of enrolled students, daily check-in counts, and pending dues.",
        "System Log Feed: Chronological logging of all data writes, login details, and payroll approvals.",
        "Quick Command Keys: Direct shortcut triggers to add members, record attendance, and launch backups."
    ]
    add_bullet_list(slide2, bullets2, 0.6, 4.05, 5.2, 3.0, font_size=12.5)

    # Screenshot on Right
    add_screenshot(slide2, images["dashboard"], 6.0, 1.4, 6.7, 5.1)

    # ==========================================
    # SLIDE 3: Faculty Attendance (Light Theme)
    # ==========================================
    slide3 = prs.slides.add_slide(slide_layout)
    apply_solid_bg(slide3, COLOR_BG_LIGHT)
    add_title(slide3, "Real-Time Staff & Faculty Attendance Sheets", COLOR_PRIMARY)

    # Content on Left
    add_card(slide3, "Paperless Check-In Records", 
             "Empower your staff clerks to log daily faculty check-ins and check-out logs digitally. Replacing register notebooks reduces entry disputes and saves accounting hours.",
             0.6, 1.4, 5.2, 2.5, badge_text="Staff Attendance", badge_color=COLOR_ACCENT_PURPLE)

    bullets3 = [
        "Flexible Status Options: Mark records as Present, Late, Half-Day, Absent, or Leave.",
        "Check-In & Out Stamps: Log exact teaching hours for audit transparency.",
        "Remarks Column: Add classroom notes, replacement remarks, or explanation details.",
        "Immediate Local Cache: Saves automatically and triggers real-time cloud backup sync."
    ]
    add_bullet_list(slide3, bullets3, 0.6, 4.05, 5.2, 3.0, font_size=12.5)

    # Screenshot on Right
    add_screenshot(slide3, images["attendance"], 6.0, 1.4, 6.7, 5.1)

    # ==========================================
    # SLIDE 4: Payroll Hub (Light Theme)
    # ==========================================
    slide4 = prs.slides.add_slide(slide_layout)
    apply_solid_bg(slide4, COLOR_BG_LIGHT)
    add_title(slide4, "One-Click Automated Payroll Hub", COLOR_PRIMARY)

    # Content on Left
    add_card(slide4, "Zero-Error Salary Computations", 
             "Payroll Hub calculates base salaries, adds dynamic allowances, and automatically subtracts leave deductions based on staff attendance logs.",
             0.6, 1.4, 5.2, 2.5, badge_text="Automated Payroll", badge_color=COLOR_ACCENT_GREEN)

    bullets4 = [
        "Leave Deductions: Integrated logic checks absences and deducts pay accurately.",
        "Allowances & Bonuses: Input custom allowances or bonuses per month.",
        "Printable Salary Slips: Generate structured printable payslips for faculty and clerks.",
        "Roll Sheet Summary: View monthly net payout metrics and approval status."
    ]
    add_bullet_list(slide4, bullets4, 0.6, 4.05, 5.2, 3.0, font_size=12.5)

    # Screenshot on Right
    add_screenshot(slide4, images["payroll"], 6.0, 1.4, 6.7, 5.1)

    # ==========================================
    # SLIDE 5: Student Enrollment Hub (Light Theme)
    # ==========================================
    slide5 = prs.slides.add_slide(slide_layout)
    apply_solid_bg(slide5, COLOR_BG_LIGHT)
    add_title(slide5, "Course Fee Master & Enrollments", COLOR_PRIMARY)

    # Content on Left
    add_card(slide5, "Revenue Control & Enrollments", 
             "Manage student registration details, parent contact records, course pricing master, discounts, and payment collections inside a clean form.",
             0.6, 1.4, 5.2, 2.5, badge_text="Student Hub", badge_color=COLOR_SECONDARY)

    bullets5 = [
        "Course Fee Master: Lock standard course fee structures to prevent unauthorized discount overrides.",
        "Detailed Profile forms: Log student name, branches, phone numbers, and parent details.",
        "Dynamic Discount Calculus: Input discount percentages or values and view Net Fees computed.",
        "Chronological Ledger: Track initial amount received and balance due for billing audits."
    ]
    add_bullet_list(slide5, bullets5, 0.6, 4.05, 5.2, 3.0, font_size=12.5)

    # Screenshot on Right
    add_screenshot(slide5, images["enrollment"], 6.0, 1.4, 6.7, 5.1)

    # ==========================================
    # SLIDE 6: Profit & Loss Ledger (Light Theme)
    # ==========================================
    slide6 = prs.slides.add_slide(slide_layout)
    apply_solid_bg(slide6, COLOR_BG_LIGHT)
    add_title(slide6, "Accounts & P&L Statement Ledgers", COLOR_PRIMARY)

    # Content on Left
    add_card(slide6, "Institute Financial Clarity", 
             "Record utility bills and operating expenses. Automated monthly accounting compares fee collections against payroll and bills to show net margins.",
             0.6, 1.4, 5.2, 2.5, badge_text="Ledger & P&L", badge_color=COLOR_ACCENT_GREEN)

    bullets6 = [
        "Bills Logging: Save utility bills (electricity, water) and other operational expenses.",
        "Auto-Computed Salaries: Automatically pulls total net payout values from Payroll Hub.",
        "Breakdown Statements: Displays itemized revenue, operating expenses, and net profit.",
        "Print P&L Reports: Export and print monthly accounting statements instantly."
    ]
    add_bullet_list(slide6, bullets6, 0.6, 4.05, 5.2, 3.0, font_size=12.5)

    # Screenshot on Right
    add_screenshot(slide6, images["ledger"], 6.0, 1.4, 6.7, 5.1)

    # ==========================================
    # SLIDE 7: Summary of Customer Benefits (Dark Theme)
    # ==========================================
    slide7 = prs.slides.add_slide(slide_layout)
    apply_solid_bg(slide7, COLOR_BG_DARK)
    add_title(slide7, "Operational ROI & Business Value", COLOR_SECONDARY)

    add_card(slide7, "Zero Spreadsheet Stress", 
             "Unifying staff rosters, payroll variables, course pricing master, student logs, and accounts into a single database removes manual data discrepancies.", 
             1.0, 1.8, 3.5, 4.5, bg_color=RGBColor(30, 41, 59), is_dark=True)

    add_card(slide7, "Delight Your Staff & Faculty", 
             "Transparent, prompt attendance check-ins, automated leave counts, and detailed printable payslip vouchers establish highly professional administration.", 
             4.9, 1.8, 3.5, 4.5, bg_color=RGBColor(30, 41, 59), is_dark=True)

    add_card(slide7, "Secure Branded Database", 
             "Run the software under your custom institute logo, address headers, theme color, and name. Service worker caching provides lightning-fast loading and offline attendance support.", 
             8.8, 1.8, 3.5, 4.5, bg_color=RGBColor(30, 41, 59), is_dark=True)

    # ==========================================
    # SLIDE 8: Conclusion & Call to Action (Dark Theme)
    # ==========================================
    slide8 = prs.slides.add_slide(slide_layout)
    apply_solid_bg(slide8, COLOR_BG_DARK)

    # Top band
    accent2 = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(0.15))
    accent2.fill.solid()
    accent2.fill.fore_color.rgb = COLOR_PRIMARY
    accent2.line.fill.background()

    cta_box = slide8.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.33), Inches(4.0))
    tf8 = cta_box.text_frame
    tf8.word_wrap = True
    
    p = tf8.paragraphs[0]
    p.text = "Upgrade Your Coaching Center Today"
    p.font.name = "Segoe UI"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLOR_SECONDARY
    p.space_after = Pt(12)
    
    p2 = tf8.add_paragraph()
    p2.text = "Transform your administrative efficiency, protect your revenues from fee leaks, and provide your office with automated bookkeeping in a sleek branded portal."
    p2.font.name = "Segoe UI"
    p2.font.size = Pt(18)
    p2.font.color.rgb = COLOR_TEXT_WHITE
    p2.space_after = Pt(24)

    p3 = tf8.add_paragraph()
    p3.text = "Contact the System Administrator to provision your branded portal and PWA app installer."
    p3.font.name = "Segoe UI"
    p3.font.size = Pt(16)
    p3.font.bold = True
    p3.font.color.rgb = COLOR_ACCENT_GREEN

    # Save presentation
    filepath = "C:/Users/deepe/.gemini/antigravity/scratch/Smart_Office_Marketing_Presentation.pptx"
    prs.save(filepath)
    print(f"Presentation created successfully at: {os.path.abspath(filepath)}")
    
    # Copy to Desktop
    import shutil
    desktop_path = "C:/Users/deepe/Desktop/Smart_Office_Marketing_Presentation.pptx"
    try:
        shutil.copy(filepath, desktop_path)
        print(f"Copied presentation to Desktop: {desktop_path}")
    except Exception as e:
        print(f"Failed to copy to Desktop: {e}")

if __name__ == "__main__":
    create_deck_with_images()
